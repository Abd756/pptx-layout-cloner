"""
Phase 3 — PPTX Reconstructor
Clones the original PPTX and surgically replaces only <a:t> text nodes
using content_mapping.json. All formatting XML is left completely untouched.
"""

import json
import os
import re
import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation

# DrawingML namespace
A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def _sanitize(text: str) -> str:
    """Strip NULL bytes and XML-illegal control characters from text."""
    if not text:
        return ""
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', str(text))


# ---------------------------------------------------------------------------
# Shape locator
# ---------------------------------------------------------------------------

def _find_shape(slide, shape_id: int):
    """Find a shape on a slide by its integer shape_id."""
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


# ---------------------------------------------------------------------------
# XML-level paragraph helpers
# ---------------------------------------------------------------------------

def _get_para_elems(shape):
    """Return all <a:p> elements from a shape's txBody."""
    return shape.text_frame._txBody.findall(f'{{{A}}}p')


def _consolidate_and_set(para_elem, new_text: str):
    """
    Core XML surgery on one <a:p>:
      1. Remove all <a:br> (line breaks)
      2. Keep only the first <a:r>, delete the rest
      3. Set that first run's <a:t> to new_text
      4. If no runs exist, create a bare <a:r><a:t> node

    Everything else in the paragraph (<a:pPr>, first run's <a:rPr>) is untouched.
    """
    # Drop line-break elements
    for br in para_elem.findall(f'{{{A}}}br'):
        para_elem.remove(br)

    runs = para_elem.findall(f'{{{A}}}r')

    if runs:
        # Keep first run's <a:rPr> intact, delete sibling runs
        first_run = runs[0]
        for run in runs[1:]:
            para_elem.remove(run)

        t = first_run.find(f'{{{A}}}t')
        if t is None:
            t = etree.SubElement(first_run, f'{{{A}}}t')
        t.text = _sanitize(new_text)

    else:
        # Paragraph had no runs (empty placeholder) — create a minimal run
        r = etree.SubElement(para_elem, f'{{{A}}}r')
        t = etree.SubElement(r, f'{{{A}}}t')
        t.text = _sanitize(new_text)


# ---------------------------------------------------------------------------
# Shape-level replacement
# ---------------------------------------------------------------------------

def _replace_single(shape, new_text: str):
    """
    For shapes with a single logical value (title, subtitle, caption).
    Puts all text in the first paragraph, blanks any extra paragraphs.
    """
    paras = _get_para_elems(shape)
    if not paras:
        return

    _consolidate_and_set(paras[0], new_text)

    # Silence any remaining paragraphs without removing them
    # (removing <a:p> nodes can corrupt the txBody structure)
    for para in paras[1:]:
        for run in para.findall(f'{{{A}}}r'):
            t = run.find(f'{{{A}}}t')
            if t is not None:
                t.text = ''


def _replace_list(shape, new_texts: list):
    """
    For shapes with multiple paragraphs (bullet lists).
    Maps new_texts[i] → paragraph[i], consolidating runs per paragraph.
    """
    paras = _get_para_elems(shape)

    for i, text in enumerate(new_texts):
        if i < len(paras):
            _consolidate_and_set(paras[i], text)
        # Phase 2 already matched counts, so i >= len(paras) shouldn't happen


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def reconstruct(pptx_path: str, content_mapping_path: str, output_path: str) -> str:
    """
    Clone *pptx_path*, replace text per *content_mapping_path*, save to *output_path*.
    Returns the output path.
    """
    # ── 1. Clone original — never modify the source ───────────────────────
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    shutil.copy2(pptx_path, output_path)
    print(f"  Cloned  : {Path(pptx_path).name} -> {Path(output_path).name}")

    # ── 2. Load content map ───────────────────────────────────────────────
    with open(content_mapping_path, 'r', encoding='utf-8') as f:
        content_map = json.load(f)

    slides_map = content_map.get('slides', {})

    # ── 3. Open the clone ─────────────────────────────────────────────────
    prs = Presentation(output_path)

    replaced = 0
    skipped  = 0

    # ── 4. Walk slides ────────────────────────────────────────────────────
    for slide_idx, slide in enumerate(prs.slides):
        slide_key    = f'slide_{slide_idx}'
        shape_map    = slides_map.get(slide_key, {})

        if not shape_map:
            print(f"  Slide {slide_idx} -- nothing to replace")
            continue

        print(f"  Slide {slide_idx} -- {len(shape_map)} shape(s) ...", end=' ', flush=True)

        for shape_id_str, new_content in shape_map.items():
            shape = _find_shape(slide, int(shape_id_str))

            if shape is None:
                print(f"\n    [WARN] shape_id {shape_id_str} not found on slide {slide_idx}")
                skipped += 1
                continue

            if not shape.has_text_frame:
                skipped += 1
                continue

            if isinstance(new_content, list):
                _replace_list(shape, new_content)
            else:
                _replace_single(shape, str(new_content))

            replaced += 1

        print('done')

    # ── 5. Save ───────────────────────────────────────────────────────────
    prs.save(output_path)

    print(f"\n[Phase 3] Reconstruction complete")
    print(f"  Shapes replaced : {replaced}")
    print(f"  Shapes skipped  : {skipped}")
    print(f"  Output          : {output_path}")

    return output_path


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser(description='Phase 3 -- Reconstruct PPTX with new content')
    parser.add_argument('--input',       required=True, help='Path to original .pptx')
    parser.add_argument('--content-map', required=True, help='Path to content_mapping.json')
    parser.add_argument('--output',      required=True, help='Path for output .pptx')
    args = parser.parse_args()

    reconstruct(args.input, args.content_map, args.output)
