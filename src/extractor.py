"""
Phase 1 — PPTX Layout Extractor
Parses any .pptx file and produces a layout.json capturing every shape's
position, styling, placeholder type, and text content.
"""

import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


# ---------------------------------------------------------------------------
# Unit helpers
# ---------------------------------------------------------------------------

def _emu_to_pt(emu):
    if emu is None:
        return None
    return round(emu / 12700, 2)


def _emu_to_pct(emu, total):
    if emu is None or not total:
        return None
    return round((emu / total) * 100, 4)


def _alignment_name(alignment):
    _map = {
        PP_ALIGN.LEFT:       "LEFT",
        PP_ALIGN.CENTER:     "CENTER",
        PP_ALIGN.RIGHT:      "RIGHT",
        PP_ALIGN.JUSTIFY:    "JUSTIFY",
        PP_ALIGN.DISTRIBUTE: "DISTRIBUTE",
    }
    return _map.get(alignment) if alignment is not None else None


def _placeholder_type_name(ph_type):
    if ph_type is None:
        return None
    try:
        return PP_PLACEHOLDER(ph_type).name
    except Exception:
        return str(ph_type)


def _safe_color_hex(color_obj):
    """Return '#RRGGBB' from a pptx ColorFormat, or None if unavailable."""
    if color_obj is None:
        return None
    try:
        if color_obj.type is not None:
            return f"#{color_obj.rgb}"
    except Exception:
        pass
    return None


# ---------------------------------------------------------------------------
# Font resolution (run → paragraph → inherited)
# ---------------------------------------------------------------------------

def _resolve_font(run, para):
    """
    Walk the font inheritance chain (run → paragraph) and return a dict
    with the most-specific concrete value found for each property.
    """
    props = {
        "font_name":    None,
        "font_size_pt": None,
        "bold":         None,
        "italic":       None,
        "underline":    None,
        "color_hex":    None,
    }

    sources = [run.font]
    try:
        sources.append(para.font)
    except AttributeError:
        pass

    for font in sources:
        if props["font_name"] is None:
            props["font_name"] = font.name or None
        if props["font_size_pt"] is None and font.size:
            props["font_size_pt"] = round(font.size.pt, 1)
        if props["bold"] is None and font.bold is not None:
            props["bold"] = font.bold
        if props["italic"] is None and font.italic is not None:
            props["italic"] = font.italic
        if props["underline"] is None and font.underline is not None:
            props["underline"] = font.underline
        if props["color_hex"] is None:
            props["color_hex"] = _safe_color_hex(font.color)

    return props


# ---------------------------------------------------------------------------
# Text-frame extraction
# ---------------------------------------------------------------------------

def _extract_run(run, para):
    font = _resolve_font(run, para)
    return {"text": run.text, **font}


def _extract_paragraph(para):
    runs = [_extract_run(r, para) for r in para.runs]

    space_before = space_after = None
    try:
        if para.space_before:
            space_before = round(para.space_before.pt, 1)
    except Exception:
        pass
    try:
        if para.space_after:
            space_after = round(para.space_after.pt, 1)
    except Exception:
        pass

    return {
        "full_text":       para.text,
        "alignment":       _alignment_name(para.alignment),
        "indent_level":    para.level,
        "space_before_pt": space_before,
        "space_after_pt":  space_after,
        "runs":            runs,
    }


def _extract_text_frame(tf):
    paragraphs = [_extract_paragraph(p) for p in tf.paragraphs]
    return {
        "full_text":       tf.text,
        "word_wrap":       tf.word_wrap,
        "auto_size":       str(tf.auto_size) if tf.auto_size else None,
        "paragraph_count": len(paragraphs),
        "paragraphs":      paragraphs,
    }


# ---------------------------------------------------------------------------
# Fill extraction
# ---------------------------------------------------------------------------

def _extract_fill(shape):
    try:
        fill = shape.fill
        fill_type = fill.type
        name = str(fill_type).split(".")[-1] if fill_type else "NONE"
        info = {"type": name}
        if name == "SOLID":
            try:
                info["color_hex"] = f"#{fill.fore_color.rgb}"
            except Exception:
                pass
        return info
    except Exception:
        return {"type": "unknown"}


# ---------------------------------------------------------------------------
# Shape extraction (with group recursion)
# ---------------------------------------------------------------------------

def _extract_single_shape(shape, slide_width, slide_height, z_order):
    data = {
        "shape_id":        shape.shape_id,
        "shape_name":      shape.name,
        "shape_type":      str(shape.shape_type).split(".")[-1] if shape.shape_type else None,
        "z_order":         z_order,
        "is_placeholder":  shape.is_placeholder,
        "placeholder_type": None,
        "placeholder_idx":  None,
        "position": {
            "left_emu":  shape.left,
            "top_emu":   shape.top,
            "width_emu": shape.width,
            "height_emu": shape.height,
            "left_pt":   _emu_to_pt(shape.left),
            "top_pt":    _emu_to_pt(shape.top),
            "width_pt":  _emu_to_pt(shape.width),
            "height_pt": _emu_to_pt(shape.height),
            "left_pct":  _emu_to_pct(shape.left,  slide_width),
            "top_pct":   _emu_to_pct(shape.top,   slide_height),
            "width_pct": _emu_to_pct(shape.width,  slide_width),
            "height_pct":_emu_to_pct(shape.height, slide_height),
        },
        "fill":            _extract_fill(shape),
        "has_text_frame":  shape.has_text_frame,
        "text_frame":      None,
    }

    if shape.is_placeholder:
        ph = shape.placeholder_format
        data["placeholder_type"] = _placeholder_type_name(ph.type)
        data["placeholder_idx"]  = ph.idx

    if shape.has_text_frame:
        data["text_frame"] = _extract_text_frame(shape.text_frame)

    return data


def _extract_shapes(shapes, slide_width, slide_height, z_start=0):
    """Recursively extract shapes, handling GROUP shapes."""
    result = []
    for z, shape in enumerate(shapes, start=z_start):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            result.append({
                "shape_id":   shape.shape_id,
                "shape_name": shape.name,
                "shape_type": "GROUP",
                "z_order":    z,
                "position": {
                    "left_emu":   shape.left,
                    "top_emu":    shape.top,
                    "width_emu":  shape.width,
                    "height_emu": shape.height,
                },
                "children": _extract_shapes(shape.shapes, slide_width, slide_height),
            })
        else:
            result.append(_extract_single_shape(shape, slide_width, slide_height, z))
    return result


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def extract_layout(pptx_path: str, output_path: str) -> dict:
    """
    Extract layout from *pptx_path* and write JSON to *output_path*.
    Returns the layout dict.
    """
    prs = Presentation(pptx_path)
    slide_width  = prs.slide_width
    slide_height = prs.slide_height

    layout = {
        "source_file":      Path(pptx_path).name,
        "slide_width_emu":  slide_width,
        "slide_height_emu": slide_height,
        "slide_width_pt":   _emu_to_pt(slide_width),
        "slide_height_pt":  _emu_to_pt(slide_height),
        "slide_count":      len(prs.slides),
        "slides":           [],
    }

    for idx, slide in enumerate(prs.slides):
        layout["slides"].append({
            "slide_index":  idx,
            "slide_id":     slide.slide_id,
            "layout_name":  slide.slide_layout.name,
            "shape_count":  len(slide.shapes),
            "shapes":       _extract_shapes(slide.shapes, slide_width, slide_height),
        })

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(layout, f, indent=2, ensure_ascii=False)

    # ── summary ──────────────────────────────────────────────────────────
    print(f"\n[Phase 1] Extraction complete")
    print(f"  Source  : {Path(pptx_path).name}")
    print(f"  Output  : {output_path}")
    print(f"  Slides  : {layout['slide_count']}")
    print(f"  Canvas  : {layout['slide_width_pt']} x {layout['slide_height_pt']} pt")
    for s in layout["slides"]:
        texts = [
            sh["text_frame"]["full_text"]
            for sh in s["shapes"]
            if sh.get("text_frame") and sh["text_frame"].get("full_text")
        ]
        print(f"  Slide {s['slide_index']} [{s['layout_name']}] — "
              f"{s['shape_count']} shapes, {len(texts)} text frames")

    return layout


# ---------------------------------------------------------------------------
# CLI usage
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Phase 1 — Extract PPTX layout to JSON")
    parser.add_argument("--input",  required=True,                   help="Path to input .pptx")
    parser.add_argument("--output", default="output/layout.json",    help="Output JSON path")
    args = parser.parse_args()

    extract_layout(args.input, args.output)
