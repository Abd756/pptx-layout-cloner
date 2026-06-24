"""
Phase 4 — Validator
Compares the original and output PPTX across layout, typography, and content
checks. Produces a validation_report.json as the proof document.
"""

import json
import os
from pathlib import Path

from lxml import etree
from pptx import Presentation

A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


# ---------------------------------------------------------------------------
# XML helpers
# ---------------------------------------------------------------------------

def _first_run_font_size(shape):
    if not shape.has_text_frame:
        return None
    try:
        paras = shape.text_frame._txBody.findall(f'{{{A}}}p')
        for para in paras:
            runs = para.findall(f'{{{A}}}r')
            for run in runs:
                rpr = run.find(f'{{{A}}}rPr')
                if rpr is not None:
                    sz = rpr.get('sz')
                    if sz:
                        return round(int(sz) / 100, 1)
    except Exception:
        pass
    return None


def _first_run_font_name(shape):
    if not shape.has_text_frame:
        return None
    try:
        paras = shape.text_frame._txBody.findall(f'{{{A}}}p')
        for para in paras:
            runs = para.findall(f'{{{A}}}r')
            for run in runs:
                rpr = run.find(f'{{{A}}}rPr')
                if rpr is not None:
                    latin = rpr.find(f'{{{A}}}latin')
                    if latin is not None:
                        return latin.get('typeface')
    except Exception:
        pass
    return None


def _first_run_color(shape):
    if not shape.has_text_frame:
        return None
    try:
        paras = shape.text_frame._txBody.findall(f'{{{A}}}p')
        for para in paras:
            runs = para.findall(f'{{{A}}}r')
            for run in runs:
                rpr = run.find(f'{{{A}}}rPr')
                if rpr is not None:
                    solid = rpr.find(f'{{{A}}}solidFill')
                    if solid is not None:
                        srgb = solid.find(f'{{{A}}}srgbClr')
                        if srgb is not None:
                            return f"#{srgb.get('val', '').upper()}"
    except Exception:
        pass
    return None


def _shape_text(shape):
    if not shape.has_text_frame:
        return ''
    return shape.text_frame.text.strip()


# ---------------------------------------------------------------------------
# Check builder
# ---------------------------------------------------------------------------

def _pass(name, **kwargs):
    return {'check': name, 'status': 'PASS', **kwargs}

def _fail(name, **kwargs):
    return {'check': name, 'status': 'FAIL', **kwargs}

def _warn(name, **kwargs):
    return {'check': name, 'status': 'WARN', **kwargs}


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def validate(original_path: str, output_path: str, report_path: str) -> dict:
    """
    Compare *original_path* vs *output_path* and write *report_path*.
    Returns the report dict.
    """
    orig = Presentation(original_path)
    out  = Presentation(output_path)

    checks = []

    # ── Global checks ─────────────────────────────────────────────────────

    # Slide count
    if len(orig.slides) == len(out.slides):
        checks.append(_pass('slide_count', value=len(orig.slides)))
    else:
        checks.append(_fail('slide_count',
                            original=len(orig.slides), output=len(out.slides)))

    # Canvas size
    if orig.slide_width == out.slide_width and orig.slide_height == out.slide_height:
        checks.append(_pass('canvas_size',
                            width_emu=orig.slide_width,
                            height_emu=orig.slide_height))
    else:
        checks.append(_fail('canvas_size',
                            original=f"{orig.slide_width}x{orig.slide_height}",
                            output=f"{out.slide_width}x{out.slide_height}"))

    # ── Per-slide checks ──────────────────────────────────────────────────

    for idx, (o_slide, out_slide) in enumerate(zip(orig.slides, out.slides)):

        # Shape count
        if len(o_slide.shapes) == len(out_slide.shapes):
            checks.append(_pass(f'shape_count_slide_{idx}',
                                value=len(o_slide.shapes)))
        else:
            checks.append(_fail(f'shape_count_slide_{idx}',
                                original=len(o_slide.shapes),
                                output=len(out_slide.shapes)))

        # Build lookup for output shapes
        out_shapes = {s.shape_id: s for s in out_slide.shapes}

        for o_sh in o_slide.shapes:
            sid  = o_sh.shape_id
            label = f'slide{idx}_sh{sid}'
            out_sh = out_shapes.get(sid)

            # Shape ID present
            if out_sh is None:
                checks.append(_fail(f'shape_exists_{label}',
                                    detail=f'shape_id {sid} missing from output'))
                continue

            checks.append(_pass(f'shape_exists_{label}'))

            # Bounding box
            pos_ok = (
                o_sh.left   == out_sh.left  and
                o_sh.top    == out_sh.top   and
                o_sh.width  == out_sh.width and
                o_sh.height == out_sh.height
            )
            if pos_ok:
                checks.append(_pass(f'bbox_{label}'))
            else:
                checks.append(_fail(f'bbox_{label}',
                                    original=f"l={o_sh.left} t={o_sh.top} w={o_sh.width} h={o_sh.height}",
                                    output=f"l={out_sh.left} t={out_sh.top} w={out_sh.width} h={out_sh.height}"))

            if not o_sh.has_text_frame:
                continue

            # Font size
            o_sz   = _first_run_font_size(o_sh)
            out_sz = _first_run_font_size(out_sh)
            if o_sz == out_sz:
                checks.append(_pass(f'font_size_{label}', value_pt=o_sz))
            elif o_sz is None or out_sz is None:
                checks.append(_warn(f'font_size_{label}',
                                    detail='inherited — not set at run level',
                                    original=o_sz, output=out_sz))
            else:
                checks.append(_fail(f'font_size_{label}',
                                    original=o_sz, output=out_sz))

            # Font family
            o_fn   = _first_run_font_name(o_sh)
            out_fn = _first_run_font_name(out_sh)
            if o_fn == out_fn:
                checks.append(_pass(f'font_family_{label}', value=o_fn))
            else:
                checks.append(_fail(f'font_family_{label}',
                                    original=o_fn, output=out_fn))

            # Font color
            o_col   = _first_run_color(o_sh)
            out_col = _first_run_color(out_sh)
            if o_col == out_col:
                checks.append(_pass(f'font_color_{label}', value=o_col))
            else:
                checks.append(_fail(f'font_color_{label}',
                                    original=o_col, output=out_col))

            # Text content
            o_text   = _shape_text(o_sh)
            out_text = _shape_text(out_sh)

            if not o_text:
                # Original was empty — decorative shape, skip text check
                checks.append(_pass(f'text_empty_ok_{label}',
                                    detail='original was empty, decorative shape'))
            elif not out_text:
                checks.append(_fail(f'text_not_empty_{label}',
                                    detail='output text is empty but original was not'))
            elif o_text != out_text:
                checks.append(_pass(f'text_replaced_{label}',
                                    original=o_text[:60],
                                    output=out_text[:60]))
            else:
                checks.append(_warn(f'text_unchanged_{label}',
                                    detail='output text same as original — LLM may have matched',
                                    value=o_text[:60]))

    # ── Summary ───────────────────────────────────────────────────────────

    passed   = sum(1 for c in checks if c['status'] == 'PASS')
    failed   = sum(1 for c in checks if c['status'] == 'FAIL')
    warnings = sum(1 for c in checks if c['status'] == 'WARN')

    report = {
        'source_original': Path(original_path).name,
        'source_output':   Path(output_path).name,
        'summary': {
            'total_checks': len(checks),
            'passed':        passed,
            'failed':        failed,
            'warnings':      warnings,
        },
        'checks': checks,
    }

    os.makedirs(os.path.dirname(os.path.abspath(report_path)), exist_ok=True)
    with open(report_path, 'w', encoding='utf-8') as f:
        json.dump(report, f, indent=2, ensure_ascii=False)

    # ── Console summary ───────────────────────────────────────────────────
    print(f"\n[Phase 4] Validation complete")
    print(f"  Total checks : {len(checks)}")
    print(f"  Passed       : {passed}")
    print(f"  Warnings     : {warnings}  (inherited styles — not failures)")
    print(f"  Failed       : {failed}")
    print(f"  Report       : {report_path}")

    if failed:
        print("\n  Failed checks:")
        for c in checks:
            if c['status'] == 'FAIL':
                print(f"    - {c['check']}: {c}")

    return report


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser(description='Phase 4 -- Validate original vs output PPTX')
    parser.add_argument('--original', required=True, help='Path to original .pptx')
    parser.add_argument('--output',   required=True, help='Path to output .pptx')
    parser.add_argument('--report',   required=True, help='Path for validation_report.json')
    args = parser.parse_args()

    validate(args.original, args.output, args.report)
